"""
PPTX Tool — Generate and read PowerPoint presentations.

Tools provided:
    create_pptx  — Generate a branded slide deck from structured content
    read_pptx    — Extract text/content from an uploaded PPTX file

Uses python-pptx for slide generation. Mezzofy brand colors:
    Orange: #f97316  |  Black: #000000  |  White: #ffffff

Slide types supported: title, content, two_column, table, thank_you.
Output files saved to the configured artifact storage directory.
"""

import logging
import uuid
from pathlib import Path
from typing import Optional

from app.tools.base_tool import BaseTool

logger = logging.getLogger("mezzofy.tools.pptx")

# Mezzofy brand colors as RGB tuples
_ORANGE = (0xF9, 0x73, 0x16)   # #f97316
_BLACK = (0x00, 0x00, 0x00)    # #000000
_WHITE = (0xFF, 0xFF, 0xFF)    # #ffffff
_LIGHT_GRAY = (0xF3, 0xF4, 0xF6)  # slide background


def _get_artifact_dir(config: dict) -> Path:
    base = config.get("storage", {}).get("local_path", "/data/artifacts")
    path = Path(base) / "pptx"
    path.mkdir(parents=True, exist_ok=True)
    return path


class PPTXOps(BaseTool):
    """PowerPoint presentation generation and reading."""

    def __init__(self, config: dict):
        super().__init__(config)
        self._artifact_dir = _get_artifact_dir(config)

    def get_tools(self) -> list[dict]:
        return [
            {
                "name": "create_pptx",
                "description": (
                    "Generate a Mezzofy-branded PowerPoint presentation from structured content. "
                    "Each slide is defined by its type and content. Supports title, content, "
                    "two-column, table, and closing slides. Returns the file path."
                ),
                "parameters": {
                    "type": "object",
                    "properties": {
                        "title": {
                            "type": "string",
                            "description": "Presentation title (shown on the cover slide).",
                        },
                        "subtitle": {
                            "type": "string",
                            "description": "Subtitle or date shown on the cover slide.",
                        },
                        "slides": {
                            "type": "array",
                            "description": "List of slide definitions.",
                            "items": {
                                "type": "object",
                                "properties": {
                                    "type": {
                                        "type": "string",
                                        "description": "Slide type.",
                                        "enum": ["title", "content", "two_column", "table", "thank_you"],
                                    },
                                    "heading": {
                                        "type": "string",
                                        "description": "Slide heading/title.",
                                    },
                                    "body": {
                                        "type": "string",
                                        "description": (
                                            "Main content. For 'content' slides: plain text or "
                                            "bullet points (one per line). For 'two_column': "
                                            "left column content."
                                        ),
                                    },
                                    "right_body": {
                                        "type": "string",
                                        "description": "Right column content (for 'two_column' slides only).",
                                    },
                                    "table_data": {
                                        "type": "array",
                                        "description": (
                                            "Table rows for 'table' slides. First row is treated as header. "
                                            "Each row is a list of cell strings."
                                        ),
                                        "items": {
                                            "type": "array",
                                            "items": {"type": "string"},
                                        },
                                    },
                                    "notes": {
                                        "type": "string",
                                        "description": "Speaker notes for this slide.",
                                    },
                                },
                                "required": ["type"],
                            },
                        },
                        "filename": {
                            "type": "string",
                            "description": "Output filename (without extension). Auto-generated if omitted.",
                        },
                    },
                    "required": ["title", "slides"],
                },
                "handler": self._create_pptx,
            },
            {
                "name": "read_pptx",
                "description": (
                    "Extract text content from an uploaded PowerPoint file. "
                    "Returns slide text, headings, and speaker notes."
                ),
                "parameters": {
                    "type": "object",
                    "properties": {
                        "file_path": {
                            "type": "string",
                            "description": "Absolute path to the PPTX file to read.",
                        },
                    },
                    "required": ["file_path"],
                },
                "handler": self._read_pptx,
            },
        ]

    def _rgb(self, rgb_tuple: tuple) -> "RGBColor":
        from pptx.util import Pt
        from pptx.dml.color import RGBColor
        return RGBColor(*rgb_tuple)

    async def _create_pptx(
        self,
        title: str,
        slides: list[dict],
        subtitle: Optional[str] = None,
        filename: Optional[str] = None,
    ) -> dict:
        """Generate a branded PPTX presentation."""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
        except ImportError:
            return self._err("python-pptx is not installed. Run: pip install python-pptx")

        if not filename:
            safe = "".join(c if c.isalnum() or c in "-_ " else "" for c in title)
            safe = safe.replace(" ", "_")[:40]
            filename = f"{safe}_{uuid.uuid4().hex[:8]}"

        output_path = self._artifact_dir / f"{filename}.pptx"

        prs = Presentation()
        # Widescreen 16:9
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        slide_layouts = prs.slide_layouts

        # Slide layout indices: 0=Title Slide, 1=Title+Content, 5=Blank
        def _set_text(tf, text: str, size: int = 18, bold: bool = False,
                      color: tuple = _BLACK, alignment=PP_ALIGN.LEFT) -> None:
            tf.word_wrap = True
            tf.text = ""
            para = tf.paragraphs[0]
            para.alignment = alignment
            run = para.add_run()
            run.text = text
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = RGBColor(*color)

        def _add_orange_bar(slide, top_inches: float = 0.0, height_inches: float = 0.12) -> None:
            """Add Mezzofy orange accent bar at top of slide."""
            bar = slide.shapes.add_shape(
                1,  # MSO_SHAPE_TYPE.RECTANGLE
                Inches(0), Inches(top_inches),
                prs.slide_width, Inches(height_inches),
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = RGBColor(*_ORANGE)
            bar.line.fill.background()

        # ── COVER SLIDE ───────────────────────────────────────────────────
        cover_layout = slide_layouts[6]  # Blank
        cover_slide = prs.slides.add_slide(cover_layout)

        # Background
        bg = cover_slide.background.fill
        bg.solid()
        bg.fore_color.rgb = RGBColor(*_BLACK)

        # Top orange bar
        _add_orange_bar(cover_slide, top_inches=0.0, height_inches=0.15)

        # Brand name
        brand_box = cover_slide.shapes.add_textbox(
            Inches(0.6), Inches(0.25), Inches(4), Inches(0.6)
        )
        _set_text(brand_box.text_frame, "Mezzofy", 22, bold=True, color=_ORANGE)

        # Title
        title_box = cover_slide.shapes.add_textbox(
            Inches(0.6), Inches(2.5), Inches(12), Inches(1.8)
        )
        _set_text(title_box.text_frame, title, 40, bold=True, color=_WHITE)

        # Subtitle
        if subtitle:
            sub_box = cover_slide.shapes.add_textbox(
                Inches(0.6), Inches(4.4), Inches(12), Inches(0.8)
            )
            _set_text(sub_box.text_frame, subtitle, 20, color=_ORANGE)

        # Bottom bar
        _add_orange_bar(cover_slide, top_inches=7.35, height_inches=0.15)

        # ── CONTENT SLIDES ─────────────────────────────────────────────────
        for slide_def in slides:
            slide_type = slide_def.get("type", "content")
            heading = slide_def.get("heading", "")
            body = slide_def.get("body", "")
            notes_text = slide_def.get("notes", "")

            if slide_type == "thank_you":
                slide = prs.slides.add_slide(slide_layouts[6])  # Blank
                _bg = slide.background.fill
                _bg.solid()
                _bg.fore_color.rgb = RGBColor(*_BLACK)
                _add_orange_bar(slide, 0.0, 0.15)
                _add_orange_bar(slide, 7.35, 0.15)

                ty_box = slide.shapes.add_textbox(
                    Inches(1), Inches(2.5), Inches(11.3), Inches(2)
                )
                _set_text(
                    ty_box.text_frame,
                    heading or "Thank You",
                    48, bold=True, color=_WHITE,
                    alignment=PP_ALIGN.CENTER,
                )
                if body:
                    sub_box = slide.shapes.add_textbox(
                        Inches(1), Inches(4.8), Inches(11.3), Inches(1)
                    )
                    _set_text(
                        sub_box.text_frame, body, 18,
                        color=_ORANGE, alignment=PP_ALIGN.CENTER,
                    )

            elif slide_type == "table" and slide_def.get("table_data"):
                slide = prs.slides.add_slide(slide_layouts[6])
                _add_orange_bar(slide, 0.0, 0.08)

                # Heading
                h_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.7))
                _set_text(h_box.text_frame, heading, 24, bold=True, color=_ORANGE)

                table_data = slide_def["table_data"]
                rows = len(table_data)
                cols = max(len(r) for r in table_data)

                tbl = slide.shapes.add_table(
                    rows, cols,
                    Inches(0.5), Inches(1.1),
                    Inches(12.3), Inches(min(5.8, rows * 0.45)),
                ).table

                for r_idx, row in enumerate(table_data):
                    for c_idx, cell_text in enumerate(row):
                        cell = tbl.cell(r_idx, c_idx)
                        cell.text = cell_text
                        tf = cell.text_frame
                        tf.paragraphs[0].runs[0].font.size = Pt(11)
                        if r_idx == 0:  # Header row
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = RGBColor(*_ORANGE)
                            tf.paragraphs[0].runs[0].font.color.rgb = RGBColor(*_WHITE)
                            tf.paragraphs[0].runs[0].font.bold = True

            elif slide_type == "two_column":
                slide = prs.slides.add_slide(slide_layouts[6])
                _add_orange_bar(slide, 0.0, 0.08)

                h_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.7))
                _set_text(h_box.text_frame, heading, 24, bold=True, color=_ORANGE)

                left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(5.8), Inches(5.8))
                _set_text(left_box.text_frame, body, 14)

                right_body = slide_def.get("right_body", "")
                right_box = slide.shapes.add_textbox(Inches(7.0), Inches(1.1), Inches(5.8), Inches(5.8))
                _set_text(right_box.text_frame, right_body, 14)

                # Divider
                div = slide.shapes.add_shape(
                    1, Inches(6.6), Inches(1.1), Inches(0.03), Inches(5.8)
                )
                div.fill.solid()
                div.fill.fore_color.rgb = RGBColor(*_ORANGE)
                div.line.fill.background()

            else:  # default: content
                slide = prs.slides.add_slide(slide_layouts[6])
                _add_orange_bar(slide, 0.0, 0.08)

                h_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.7))
                _set_text(h_box.text_frame, heading, 24, bold=True, color=_ORANGE)

                body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(12.3), Inches(5.8))
                body_box.text_frame.word_wrap = True

                # Render bullet points if body has newlines
                lines = body.strip().split("\n") if body else []
                if lines:
                    tf = body_box.text_frame
                    tf.text = ""
                    for i, line in enumerate(lines):
                        if i == 0:
                            para = tf.paragraphs[0]
                        else:
                            para = tf.add_paragraph()
                        para.text = line.lstrip("•-* ")
                        para.level = 0
                        run = para.runs[0] if para.runs else para.add_run()
                        run.font.size = Pt(16)
                        run.font.color.rgb = RGBColor(*_BLACK)

            # Speaker notes
            if notes_text:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = notes_text

        prs.save(str(output_path))
        file_size = output_path.stat().st_size

        logger.info(f"Created PPTX: {output_path} ({len(slides)+1} slides)")
        return self._ok({
            "file_path": str(output_path),
            "filename": f"{filename}.pptx",
            "slide_count": len(slides) + 1,  # +1 for cover
            "size_bytes": file_size,
            "title": title,
        })

    async def _read_pptx(self, file_path: str) -> dict:
        """Extract text from a PPTX file."""
        import os
        if not os.path.exists(file_path):
            return self._err(f"File not found: {file_path}")

        try:
            from pptx import Presentation
        except ImportError:
            return self._err("python-pptx is not installed. Run: pip install python-pptx")

        try:
            prs = Presentation(file_path)
            slides = []

            for i, slide in enumerate(prs.slides):
                texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                texts.append(text)

                notes_text = ""
                if slide.has_notes_slide:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()

                slides.append({
                    "slide": i + 1,
                    "text": texts,
                    "notes": notes_text,
                })

            return self._ok({
                "file_path": file_path,
                "slide_count": len(prs.slides),
                "slides": slides,
            })

        except Exception as e:
            logger.error(f"Failed to read PPTX {file_path}: {e}")
            return self._err(str(e))
