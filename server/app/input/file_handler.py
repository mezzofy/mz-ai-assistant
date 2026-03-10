"""
File Handler — text extraction from uploaded documents.

Supported formats and extraction libraries:
  PDF   → PDFOps.read_pdf (pypdf)
  DOCX  → DocxOps.read_docx (paragraphs + tables)
  PPTX  → PPTXOps.read_pptx (slide text + speaker notes)
  CSV   → CSVOps.read_csv (up to 500 rows + numeric summary)
  XLSX  → pandas
  DOC   → python-docx (legacy Word 97-2003, best-effort)
  PPT   → python-pptx (legacy PowerPoint 97-2003, best-effort)
  TXT   → plain UTF-8 read
"""

import logging
import os
import tempfile
from typing import Optional

logger = logging.getLogger("mezzofy.input.file")

# Truncate extracted text at this many chars to avoid overwhelming LLM context
_MAX_EXTRACTED_CHARS = 6000


async def handle_file(
    task: dict,
    file_bytes: bytes,
    filename: str,
) -> dict:
    """
    Extract text from an uploaded document file.

    Args:
        task:       Task dict with _config and optional message.
        file_bytes: Raw document bytes.
        filename:   Original filename (used for extension detection).

    Returns:
        Task dict enriched with extracted_text from the document content.
    """
    config = task.get("_config", {})
    ext = os.path.splitext(filename)[1].lower()

    tmp_path: Optional[str] = None
    try:
        with tempfile.NamedTemporaryFile(suffix=ext or ".bin", delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name

        extracted = await _extract_by_extension(ext, tmp_path, config)
        extracted = extracted[:_MAX_EXTRACTED_CHARS]

        parts = []
        if extracted:
            parts.append(
                f"[Document content from '{filename}':\n{extracted}]"
            )
        user_msg = (task.get("message") or "").strip()
        if user_msg:
            parts.append(user_msg)

        return {
            **task,
            "input_type": "file",
            "extracted_text": (
                "\n\n".join(parts)
                if parts
                else f"[File '{filename}' uploaded — no text extracted]"
            ),
            "media_content": {
                "filename": filename,
                "extension": ext,
                "extracted_chars": len(extracted),
            },
            "input_summary": (
                f"File: {filename} ({len(extracted):,} chars extracted)"
            ),
        }

    finally:
        if tmp_path:
            try:
                os.unlink(tmp_path)
            except Exception:
                pass


async def _extract_by_extension(ext: str, file_path: str, config: dict) -> str:
    """Route to the correct extraction method by file extension."""

    if ext == ".pdf":
        try:
            from app.tools.document.pdf_ops import PDFOps
            ops = PDFOps(config)
            result = await ops.execute("read_pdf", file_path=file_path)
            if not result.get("success"):
                return ""
            data = result.get("output", {})
            pages = data.get("pages", []) if isinstance(data, dict) else []
            return "\n\n".join(p.get("text", "") for p in pages if p.get("text"))
        except Exception as e:
            logger.warning(f"PDF extraction failed: {e}")
            return ""

    if ext == ".docx":
        try:
            from app.tools.document.docx_ops import DocxOps
            result = await DocxOps(config).execute("read_docx", file_path=file_path)
            data = result.get("output", {}) if result.get("success") else {}
            parts = [p["text"] for p in data.get("paragraphs", []) if p.get("text")]
            for table in data.get("tables", []):
                for row in table.get("rows", []):
                    parts.append("  |  ".join(cell for cell in row if cell))
            return "\n".join(parts)
        except Exception as e:
            logger.warning(f"DOCX extraction failed: {e}")
            return ""

    if ext == ".pptx":
        try:
            from app.tools.document.pptx_ops import PPTXOps
            result = await PPTXOps(config).execute("read_pptx", file_path=file_path)
            data = result.get("output", {}) if result.get("success") else {}
            parts = []
            for slide in data.get("slides", []):
                parts.extend(slide.get("text", []))
                if slide.get("notes"):
                    parts.append(f"[Notes: {slide['notes']}]")
            return "\n".join(parts)
        except Exception as e:
            logger.warning(f"PPTX extraction failed: {e}")
            return ""

    if ext == ".csv":
        try:
            from app.tools.document.csv_ops import CSVOps
            result = await CSVOps(config).execute("read_csv", file_path=file_path, max_rows=500)
            data = result.get("output", {}) if result.get("success") else {}
            headers = data.get("headers", [])
            rows = data.get("rows", [])
            lines = ["  |  ".join(str(h) for h in headers)] if headers else []
            for row in rows:
                lines.append("  |  ".join(str(v) for v in row))
            summary = data.get("numeric_summary", {})
            if summary:
                lines.append("\n[Column Summary]")
                for col, stats in summary.items():
                    lines.append(f"  {col}: min={stats.get('min')}, max={stats.get('max')}, mean={stats.get('mean'):.2f}")
            return "\n".join(lines)
        except Exception as e:
            logger.warning(f"CSV extraction failed: {e}")
            return ""

    if ext == ".doc":
        try:
            import docx as python_docx
            doc = python_docx.Document(file_path)
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        except Exception:
            return (
                "[Note: This file is in legacy Word 97-2003 (.doc) format which could not be "
                "fully parsed. For best results, please re-save as .docx and re-upload.]"
            )

    if ext == ".ppt":
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            texts = [shape.text.strip() for slide in prs.slides
                     for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
            return "\n".join(texts)
        except Exception:
            return (
                "[Note: This file is in legacy PowerPoint 97-2003 (.ppt) format which could not be "
                "fully parsed. For best results, please re-save as .pptx and re-upload.]"
            )

    if ext in (".xlsx", ".xls"):
        try:
            import pandas as pd
            df = pd.read_excel(file_path, nrows=100)
            return df.to_string(index=False)
        except Exception as e:
            logger.warning(f"Excel extraction failed: {e}")
            return ""

    if ext in (".txt", ".md", ".rst"):
        try:
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                return f.read()
        except Exception as e:
            logger.warning(f"Text read failed: {e}")
            return ""

    logger.warning(f"No extractor for extension: {ext!r}")
    return ""
